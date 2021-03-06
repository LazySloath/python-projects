# -*- coding: utf-8 -*-
"""
Created on Mon Jan  8 21:47:27 2018

@author: User

Powerpoint lyrics generator that scrapes azlyrics for lyrics
to be implemented: more design choices
Added function to review singer/title

now
IP BANNED D':
still problems with whitespace
design tools
lines per slide? how to fix
"""

import tkinter, requests, bs4, docx, pptx
import tkinter.filedialog
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

class PptLyrics(tkinter.Tk):
    # Parent is its container
    def __init__(self, parent):
        tkinter.Tk.__init__(self, parent)
        self.parent = parent
        self.initialise()
        
    def initialise(self):
        self.state = 0
        self.grid()
        # Set window size
        self.geometry("410x85")
        # Set window icon
        self.wm_iconbitmap('chordsearch.ico')
        self.entryVariable = tkinter.StringVar()
        # Entry field for song titles
        self.entry = tkinter.Entry(self, textvariable = self.entryVariable)
        self.entry.grid(column = 0, row = 2, columnspan = 5, sticky = 'EW')
        self.entryVariable.set('Song title here')
        # Hook enter so that OnPressEnter occurs when enter is pressed
        self.entry.bind('<Return>', self.OnPressEnter)
        # Enter button
        enter = tkinter.Button(self, text = 'Enter', 
                                command = self.OnEnterClick)
        enter.grid(column = 3, row = 3)
        # Save button
        save = tkinter.Button(self, text = 'Save', 
                                command = self.OnSaveClick)
        save.grid(column = 4, row = 3)
        # Up button
        up = tkinter.Button(self, text = '/\\', command = self.OnUpClick)
        up.grid(column = 1, row = 3)
        # Down button
        down = tkinter.Button(self, text = '\\/', command = self.OnDownClick)
        down.grid(column = 2, row = 3)
        # Label for instructions
        self.instructionVariable = tkinter.StringVar()
        instruction = tkinter.Label(self, anchor = 'w', 
                                    fg = 'white', bg = 'black', 
                                    textvariable = self.instructionVariable)
        self.instructionVariable.set('Enter a song title '
                                    + 'to search for lyrics or press the '
                                    + '"Save" button when done.')
        instruction.grid(column = 0, row = 0, columnspan = 5, sticky = 'EW')
        # Label for status
        self.labelVariable = tkinter.StringVar()
        label = tkinter.Label(self, textvariable = self.labelVariable,
                              anchor = 'w', fg = 'white', bg = 'grey')
        label.grid(column = 0, row = 1, columnspan = 5, sticky = 'EW')
        # Don't allow grid to resize
        self.grid_columnconfigure(0, weight = 1)
        self.resizable(False, False)
        # Don't allow auto grow and shrink for window
        self.update()
        self.geometry(self.geometry())
            
    def OnEnterClick(self):
        if self.state == 0:
        # Get song title from entry field
            title = self.entryVariable.get()
            self.songlist = song_search(title)
            # No lyrics found
            if self.songlist == []:
                self.labelVariable.set('No lyrics found for ' + title + '!')
                # Lyrics found
            else:
                self.instructionVariable.set('Use the arrows to search for ' 
                                             + 'the correct version and press ' 
                                             + '"Enter" to continue.')
                self.counter = 0
                self.labelVariable.set(str(self.counter + 1) + ': '
                                       + self.songlist[self.counter][0] + ' / ' 
                                       + self.songlist[self.counter][1])
                self.state = 1
                
        elif self.state == 1:
            title = self.songlist[self.counter][1]
            lyrics = self.songlist[self.counter][2].getText()
            print(lyrics)
            lyriclist = cleanup(lyrics)
            add_lyrics(title, lyriclist)
            self.labelVariable.set('Lyrics added for ' + title + '!')
            self.instructionVariable.set('Enter a song title '
                                         + 'to search for lyrics or press the '
                                         + '"Save" button when done.')
            self.state = 0
            
        # Auto keyboard focus on entry field
        self.entry.focus_set()
        self.entry.selection_range(0, tkinter.END)
    
    def OnPressEnter(self, event):
        self.OnEnterClick()
        
    def OnSaveClick(self):
        # Get user input for doc name then exit
        self.filepath = tkinter.filedialog.asksaveasfilename()
        app.destroy()
        
    def OnUpClick(self):
        if self.state == 0:
            pass
        
        elif self.state == 1:
            # Scroll through versions
            if self.counter == 0:
                pass
            else:
                self.counter -= 1
                self.labelVariable.set(str(self.counter + 1) + ': '
                                       + self.songlist[self.counter][0] + ' / ' 
                                       + self.songlist[self.counter][1])
                
            
    def OnDownClick(self):
        if self.state == 0:
            pass
        
        elif self.state == 1:
            # Scroll through versions
            if self.counter == len(self.songlist) - 1:
                pass
            else:
                self.counter += 1
                self.labelVariable.set(str(self.counter + 1) + ': '
                                       + self.songlist[self.counter][0] + ' / ' 
                                       + self.songlist[self.counter][1])
                 

def soupify(link):
    """Create BeautifulSoup object from a link"""
    
    # Spoof user agent to prevent connection error
    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 6.0; WOW64; rv:24.0) Gecko/20100101 Firefox/24.0'}

    res = requests.get(link, headers = headers)
    res.raise_for_status
    return bs4.BeautifulSoup(res.text, 'lxml')

def song_search(title):
    """Search for song on azlyrics, return list of tuples(singer,title,lyrics)"""
    
    ugsearch = ('https://search.azlyrics.com/search.php?q=' + title)
    ugsoup = soupify(ugsearch)
    # Get list 'hits' of search results
    hits = ugsoup.select('td > a')
    singers = ugsoup.select('td > b')
    counter = 0
    songlist = []
    for hit in hits:
        href = hit.get('href')
        if 'https://www.azlyrics.com/lyrics/' in href:
            title = hit.getText()
            soup = soupify(href)
            # Find plain div element (how azlyrics stores lyrics)
            lyrics = soup.find('div', class_ = False, id = False)
            songlist.append((singers[counter].getText(), title, lyrics))
            counter += 1
            
    return songlist
    
def cleanup(lyrics):
    """Clean up lyrics and return lyrics in list form"""
    
    lyriclist = lyrics.split('\n')
    try:
        lyriclist.remove('\r')
    except:
        pass
    # Remove [verse]
    for line in lyriclist.copy():
        try:
            if line[0] == '[':
                lyriclist.remove(line)
        except IndexError:
            pass
    # Remove double ''
    removals = []
    for i in range(len(lyriclist) - 1):
        if lyriclist[i] == '' and lyriclist[i+1] == '':
            removals.append(i)
    for i in reversed(removals):
        del lyriclist[i]
    # Remove last ''
    if lyriclist[-1] == '':
        del lyriclist[-1]
    return lyriclist

def add_lyrics(title, lyriclist):
    """Add chords to the ppt"""
    title = capitalise(title)
    # Iterate over lyriclist, '' demarcates new verse
    for line in lyriclist:
        if line == '' or np == 4:
            slide = prs.slides.add_slide(titlecontent)
            slide.placeholders[0].text = title
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            np = 0
        else:
            # Select nth paragraph to insert nth line
            p = text_frame.paragraphs[np]
            # Make bullet point invisible
            emptyrun = p.add_run()
            emptyrun.text = ' '
            emptyrun.font.size = Pt(1)
            emptyrun.font.color.rgb = RGBColor(255,255,255)
            # Add line
            run = p.add_run()
            run.text = line
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(0,0,0)
            np += 1
            text_frame.add_paragraph()


def capitalise(title):
    """Capitalise each word in the title"""
    titlelist = []
    for word in title.split():
        try:
            titlelist.append(word[0].upper() + word[1:].lower())
        except IndexError:
            titlelist.append(word.upper())
    return ' '.join(titlelist)

if __name__ == "__main__":
    # Initialise PPT
    prs = pptx.Presentation()
    # Set layout
    titleslide = prs.slide_layouts[0]
    titlecontent = prs.slide_layouts[1]
    # Intro slide
    introslide = prs.slides.add_slide(titleslide)
    introslide.placeholders[0].text = 'Welcome to Cell :)'
    
    # Run GUI to let user get lyrics as required    
    app = PptLyrics(None)
    app.title('EasiLyric')
    app.mainloop()
    # Save the chordsheet
    try:
        if len(app.filepath) > 0:
            prs.save(app.filepath + '.pptx')
    except AttributeError:
        pass